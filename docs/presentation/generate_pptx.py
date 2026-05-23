"""
Generate Legal AI AR pitch deck (.pptx) matching the HTML presentation design.
Usage:  python docs/presentation/generate_pptx.py
Output: docs/presentation/legal-ai-ar-pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

ORANGE = RGBColor(0xD0, 0x4A, 0x02)
AMBER = RGBColor(0xEB, 0x8C, 0x00)
GRAY_900 = RGBColor(0x25, 0x25, 0x25)
GRAY_800 = RGBColor(0x47, 0x47, 0x47)
GRAY_500 = RGBColor(0x69, 0x69, 0x69)
GRAY_300 = RGBColor(0xD1, 0xD1, 0xD1)
GRAY_100 = RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x21, 0x81, 0x2D)
SUCCESS_BG = RGBColor(0xE9, 0xF5, 0xEB)
ORANGE_BG = RGBColor(0xFF, 0xF0, 0xEB)
ORANGE_BORDER = RGBColor(0xF5, 0xC4, 0xB0)
HERO_BG = RGBColor(0xFF, 0xFA, 0xF7)
DARK_CARD = RGBColor(0x3A, 0x3A, 0x3A)
DARK_BORDER = RGBColor(0x50, 0x50, 0x50)
LIGHT_GRAY_BG = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_GRAY_FG = RGBColor(0x55, 0x55, 0x55)

HD = "Georgia"
BD = "Arial"
N = 13
SW = Inches(13.333)
SH = Inches(7.5)


def _sl(prs, bg=WHITE):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return sl


def _tb(sl, l, t, w, h, txt, *, fnt=BD, sz=Pt(14), clr=GRAY_800,
        b=False, it=False, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = al
    r = p.add_run()
    r.text = txt
    r.font.name = fnt
    r.font.size = sz
    r.font.color.rgb = clr
    r.font.bold = b
    r.font.italic = it
    return tb


def _eye(sl, top, txt, clr=ORANGE):
    return _tb(sl, Inches(1), top, Inches(11.333), Pt(20),
               txt.upper(), sz=Pt(11), clr=clr, b=True, al=PP_ALIGN.CENTER)


def _h(sl, top, txt, sz=Pt(28), clr=GRAY_900, al=PP_ALIGN.CENTER, w=Inches(11.333)):
    left = (SW - w) // 2
    return _tb(sl, left, top, w, Inches(1.2), txt,
               fnt=HD, sz=sz, clr=clr, al=al)


def _foot(sl, n, clr=GRAY_500):
    _tb(sl, Inches(0.75), Inches(6.95), Inches(5), Inches(0.3),
        "PwC Argentina \u2014 Marzo 2026 \u2014 Confidencial", sz=Pt(9), clr=clr)
    _tb(sl, Inches(11), Inches(6.95), Inches(1.5), Inches(0.3),
        f"{n} / {N}", sz=Pt(9), clr=clr, al=PP_ALIGN.RIGHT)


def _rr(sl, l, t, w, h, *, fill=WHITE, line=GRAY_300, lw=Pt(0.75), rad=0.02):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = lw
    sh.adjustments[0] = rad
    return sh


def _acc(sl, l, t, h, clr=ORANGE):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(3), h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = clr
    sh.line.fill.background()


def _ib(sl, l, t, sym):
    s = Inches(0.38)
    _rr(sl, l, t, s, s, fill=ORANGE_BG, line=ORANGE_BG, rad=0.05)
    _acc(sl, l, t, s)
    _tb(sl, l, t - Pt(1), s, s, sym, sz=Pt(14), clr=ORANGE, al=PP_ALIGN.CENTER)


def _card(sl, l, t, w, h, *, icon=None, title="", body=""):
    _rr(sl, l, t, w, h)
    _acc(sl, l, t + Pt(4), h - Pt(8))
    y = t + Inches(0.22)
    if icon:
        _ib(sl, l + Inches(0.22), y, icon)
        y += Inches(0.5)
    if title:
        _tb(sl, l + Inches(0.22), y, w - Inches(0.44), Inches(0.25),
            title, sz=Pt(12), clr=GRAY_900, b=True)
        y += Inches(0.32)
    if body:
        _tb(sl, l + Inches(0.22), y, w - Inches(0.44),
            h - (y - t) - Inches(0.15), body, sz=Pt(11), clr=GRAY_800)


def _stat(sl, l, t, w, h, val, label, *, vc=ORANGE):
    _rr(sl, l, t, w, h)
    _tb(sl, l, t + Inches(0.25), w, Inches(0.55),
        str(val), fnt=HD, sz=Pt(36), clr=vc, al=PP_ALIGN.CENTER)
    _tb(sl, l, t + Inches(0.85), w, Inches(0.3),
        label.upper(), sz=Pt(9), clr=GRAY_500, b=True, al=PP_ALIGN.CENTER)


def _badges(sl, l, t, items, *, bg=ORANGE_BG, fg=ORANGE,
            gap=Inches(0.08), max_w=None, fsz=Pt(9)):
    x, y = l, t
    cpi = 5.0 if fsz <= Pt(8) else 5.8
    bh = Pt(20) if fsz <= Pt(8) else Pt(22)
    for txt in items:
        w = Pt(len(txt) * cpi + 24)
        if max_w and (x + w - l) > max_w:
            x = l
            y += Pt(28)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, bh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = bg
        sh.line.fill.background()
        sh.adjustments[0] = 0.5
        tf = sh.text_frame
        tf.word_wrap = False
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.name = BD
        r.font.size = fsz
        r.font.color.rgb = fg
        r.font.bold = True
        x += w + gap


def _pipe(sl, x, y, txt, *, bg=ORANGE_BG, fg=ORANGE, bd=ORANGE_BORDER):
    w, h = Inches(1.6), Inches(0.45)
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = bd
    sh.line.width = Pt(0.75)
    sh.adjustments[0] = 0.04
    _tb(sl, x, y + Pt(2), w, h, txt, sz=Pt(11), clr=fg, b=True, al=PP_ALIGN.CENTER)
    return w


def _arw(sl, x, y, length=Inches(0.5)):
    _tb(sl, x, y - Pt(4), length, Inches(0.45),
        "\u2192", sz=Pt(18), clr=GRAY_300, al=PP_ALIGN.CENTER)
    return length


# ─── SLIDES ──────────────────────────────────────────────────

def s00_cover(prs):
    sl = _sl(prs, HERO_BG)
    _tb(sl, Inches(5.1), Inches(1.7), Inches(1.1), Inches(0.4),
        "pwc", fnt=HD, sz=Pt(26), clr=ORANGE, b=True)
    sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(6.25), Inches(1.78), Pt(1), Inches(0.25))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GRAY_300
    sep.line.fill.background()
    _tb(sl, Inches(6.45), Inches(1.72), Inches(2), Inches(0.4),
        "Legal AI AR", sz=Pt(14), clr=GRAY_900)
    _eye(sl, Inches(2.35), "Plataforma de Inteligencia Artificial Legal")
    _h(sl, Inches(2.8), "Jurisprudencia argentina\nal alcance de una pregunta", sz=Pt(34))
    _tb(sl, Inches(2), Inches(4.2), Inches(9.333), Inches(0.7),
        "Ingesta automatizada, b\u00fasqueda sem\u00e1ntica y asistente jurisprudencial "
        "con IA generativa para potenciar la pr\u00e1ctica legal.",
        sz=Pt(14), clr=GRAY_500, al=PP_ALIGN.CENTER)
    _badges(sl, Inches(3.7), Inches(5.2),
            ["Azure Cloud-Native", ".NET 8 + Angular 18", "GPT-5o"])
    _foot(sl, 1)


def s01_problem(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(1.2), "El desaf\u00edo")
    _h(sl, Inches(1.55),
       "La investigaci\u00f3n jurisprudencial hoy es lenta,\nfragmentada e incompleta",
       sz=Pt(26))
    cw, ch = Inches(3.4), Inches(2.6)
    gap = Inches(0.35)
    x0 = (SW - 3 * cw - 2 * gap) // 2
    y0 = Inches(3.1)
    data = [
        ("\U0001f550", "Horas de b\u00fasqueda manual",
         "Los abogados navegan m\u00faltiples portales p\u00fablicos (CSJN, SAIJ, PJN, SCBA) "
         "sin integraci\u00f3n, copiando y pegando resultados."),
        ("\U0001f50d", "B\u00fasqueda por texto exacto",
         "Los buscadores p\u00fablicos solo soportan keywords. No entienden la intenci\u00f3n "
         "ni el contexto jur\u00eddico de la consulta."),
        ("\u26a0", "Riesgo de fallos omitidos",
         "Sin cobertura completa ni relaciones entre fallos, se pierden precedentes "
         "relevantes que pueden definir un caso."),
    ]
    for i, (ic, ti, bo) in enumerate(data):
        _card(sl, x0 + i * (cw + gap), y0, cw, ch, icon=ic, title=ti, body=bo)
    _foot(sl, 2)


def s02_solution(prs):
    sl = _sl(prs, HERO_BG)
    _eye(sl, Inches(0.7), "La soluci\u00f3n")
    _h(sl, Inches(1.0), "Legal AI AR: tu asistente jurisprudencial con IA", sz=Pt(26))

    bullets = [
        ("Ingesta autom\u00e1tica", " de fallos desde fuentes p\u00fablicas argentinas"),
        ("B\u00fasqueda sem\u00e1ntica", " en lenguaje natural sobre toda la base de conocimiento"),
        ("Chat RAG ag\u00e9ntico", " que responde citando fallos concretos con links"),
        ("Grafo de citas", " entre fallos con resoluci\u00f3n retroactiva"),
        ("Panel de administraci\u00f3n", " para controlar el pipeline de ingesta"),
        ("Guardrails de seguridad", ": moderaci\u00f3n, verificaci\u00f3n de citas, disclaimer"),
    ]
    y = Inches(2.0)
    for bold_txt, rest_txt in bullets:
        tb = sl.shapes.add_textbox(Inches(1), y, Inches(5.5), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r0 = p.add_run()
        r0.text = "\u25cf  "
        r0.font.name = BD
        r0.font.size = Pt(8)
        r0.font.color.rgb = ORANGE
        r1 = p.add_run()
        r1.text = bold_txt
        r1.font.name = BD
        r1.font.size = Pt(12)
        r1.font.color.rgb = GRAY_900
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = rest_txt
        r2.font.name = BD
        r2.font.size = Pt(12)
        r2.font.color.rgb = GRAY_800
        y += Inches(0.38)

    cx, cy, cw2, ch2 = Inches(7), Inches(2.0), Inches(5.5), Inches(4.2)
    _rr(sl, cx, cy, cw2, ch2, fill=GRAY_100, line=GRAY_100)
    _tb(sl, cx + Inches(0.3), cy + Inches(0.2), cw2 - Inches(0.6), Inches(0.25),
        "Ejemplo de consulta al chat:", sz=Pt(11), clr=GRAY_500)

    qy = cy + Inches(0.55)
    qh = Inches(0.75)
    _rr(sl, cx + Inches(0.3), qy, cw2 - Inches(0.6), qh, fill=ORANGE_BG, line=ORANGE_BG)
    _acc(sl, cx + Inches(0.3), qy, qh)
    _tb(sl, cx + Inches(0.55), qy + Inches(0.1), cw2 - Inches(0.9), qh - Inches(0.15),
        "\u00bfQu\u00e9 dijo la Corte Suprema sobre la doble imposici\u00f3n "
        "tributaria en los \u00faltimos 5 a\u00f1os?",
        sz=Pt(12), clr=GRAY_900)

    ay = qy + qh + Inches(0.15)
    _tb(sl, cx + Inches(0.3), ay, cw2 - Inches(0.6), Inches(0.2),
        "Respuesta del asistente:", sz=Pt(11), clr=GRAY_500)
    ay += Inches(0.3)
    ah = Inches(1.6)
    _rr(sl, cx + Inches(0.3), ay, cw2 - Inches(0.6), ah)

    tb = sl.shapes.add_textbox(cx + Inches(0.5), ay + Inches(0.15),
                                cw2 - Inches(1.0), ah - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "La CSJN ha abordado la doble imposici\u00f3n en varios fallos recientes. En "
    r.font.name = BD
    r.font.size = Pt(11)
    r.font.color.rgb = GRAY_800
    r2 = p.add_run()
    r2.text = "\"Molinos R\u00edo de la Plata S.A. c/ DGI\""
    r2.font.name = BD
    r2.font.size = Pt(11)
    r2.font.color.rgb = ORANGE
    r2.font.bold = True
    r3 = p.add_run()
    r3.text = " (2022) estableci\u00f3 que..."
    r3.font.name = BD
    r3.font.size = Pt(11)
    r3.font.color.rgb = GRAY_800

    _tb(sl, cx + Inches(0.5), ay + ah - Inches(0.4), cw2 - Inches(1.0), Inches(0.25),
        "\u2713 3 citas verificadas contra la base de conocimiento",
        sz=Pt(9), clr=GRAY_500)
    _foot(sl, 3)


def s03_architecture(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(0.7), "Arquitectura")
    _h(sl, Inches(1.0),
       "Pipeline inteligente de ingesta + Knowledge Base h\u00edbrida", sz=Pt(24))

    steps = ["Crawler", "Parser", "Enrichment\n(GPT-5o)", "Indexer", "Knowledge\nBase"]
    s_bg = [ORANGE_BG] * 4 + [SUCCESS_BG]
    s_fg = [ORANGE] * 4 + [SUCCESS]
    s_bd = [ORANGE_BORDER] * 4 + [RGBColor(0xB5, 0xDD, 0xB9)]
    sw = Inches(1.7)
    aw = Inches(0.45)
    total = 5 * sw + 4 * aw
    x = (SW - total) // 2
    py = Inches(2.1)
    for i, st in enumerate(steps):
        _pipe(sl, x, py, st, bg=s_bg[i], fg=s_fg[i], bd=s_bd[i])
        x += sw
        if i < 4:
            _arw(sl, x, py)
            x += aw

    cw, ch = Inches(3.4), Inches(1.8)
    gap = Inches(0.35)
    x0 = (SW - 3 * cw - 2 * gap) // 2
    y0 = Inches(3.3)
    cards = [
        ("Azure SQL",
         "Metadata estructurada, relaciones entre fallos, jueces, leyes y voces tem\u00e1ticas. "
         "Grafo de citas con CTEs recursivas."),
        ("Azure AI Search",
         "B\u00fasqueda h\u00edbrida (vector + keyword) sobre 2 \u00edndices: por fallo y por "
         "chunk de 512 tokens con overlap 50."),
        ("Azure Blob + OpenAI",
         "PDFs originales en Blob. Embeddings con text-embedding-3-large (3072 dims). "
         "Chat con gpt-5o."),
    ]
    for i, (ti, bo) in enumerate(cards):
        _card(sl, x0 + i * (cw + gap), y0, cw, ch, title=ti, body=bo)

    _badges(sl, Inches(1.5), Inches(5.5),
            ["Deduplicaci\u00f3n SHA-256", "Idempotencia", "Scale-to-zero (KEDA)",
             "DLQ + Retry", "Resoluci\u00f3n retroactiva de citas"],
            bg=LIGHT_GRAY_BG, fg=LIGHT_GRAY_FG)
    _foot(sl, 4)


def s04_agentic(prs):
    sl = _sl(prs, HERO_BG)
    _eye(sl, Inches(0.5), "Funcionalidad estrella")
    _h(sl, Inches(0.8),
       "Recorrido de una consulta por el pipeline ag\u00e9ntico", sz=Pt(24))
    _tb(sl, Inches(2), Inches(1.45), Inches(9.333), Inches(0.4),
        "Cada pregunta del abogado atraviesa 5 agentes especializados. "
        "El Executor invoca herramientas din\u00e1micamente antes de responder.",
        sz=Pt(11), clr=GRAY_500, al=PP_ALIGN.CENTER)

    nodes = [
        ("\U0001f464", "Abogado", None, WHITE, GRAY_900),
        ("\U0001f6e1", "Input\nGuardrail",
         "Filtra prompt injection,\nPII y fuera de scope", WHITE, GRAY_900),
        ("\u2b50", "Query\nEnricher",
         "Detecta intenci\u00f3n,\nextrae entidades", WHITE, GRAY_900),
        ("\u2699", "Agentic\nExecutor",
         "GPT-5o decide qu\u00e9\nherramientas invocar", ORANGE_BG, ORANGE),
        ("\u2713", "Output\nGuardrail",
         "Valida citas contra\nla DB + disclaimer", WHITE, GRAY_900),
        ("\U0001f4c4", "Finalizer",
         "Normaliza citas,\nformato final", WHITE, GRAY_900),
        ("\U0001f4ac", "Respuesta\nverificada", None, SUCCESS_BG, SUCCESS),
    ]
    nw, nh = Inches(1.35), Inches(1.3)
    aw2 = Inches(0.35)
    total = len(nodes) * nw + (len(nodes) - 1) * aw2
    nx = (SW - total) // 2
    ny = Inches(2.2)

    for i, (ic, label, desc, bg, fg) in enumerate(nodes):
        _rr(sl, nx, ny, nw, nh, fill=bg,
            line=ORANGE if fg == ORANGE else GRAY_300)
        _tb(sl, nx, ny + Inches(0.08), nw, Inches(0.3),
            ic, sz=Pt(18), al=PP_ALIGN.CENTER, clr=fg)
        _tb(sl, nx + Inches(0.05), ny + Inches(0.42), nw - Inches(0.1), Inches(0.35),
            label, sz=Pt(10), clr=fg, b=True, al=PP_ALIGN.CENTER)
        if desc:
            _tb(sl, nx + Inches(0.05), ny + Inches(0.78), nw - Inches(0.1), Inches(0.5),
                desc, sz=Pt(8), clr=GRAY_500, al=PP_ALIGN.CENTER)
        nx += nw
        if i < len(nodes) - 1:
            _tb(sl, nx, ny + Inches(0.4), aw2, Inches(0.3),
                "\u2192", sz=Pt(16), clr=GRAY_300, al=PP_ALIGN.CENTER)
            nx += aw2

    _tb(sl, Inches(2), Inches(3.8), Inches(9.333), Inches(0.25),
        "HERRAMIENTAS DISPONIBLES (FUNCTION CALLING)", sz=Pt(9),
        clr=GRAY_500, b=True, al=PP_ALIGN.CENTER)

    tools = ["Buscar fallos", "Detalle de fallo", "Citas y relaciones",
             "Buscar por ley", "Estad\u00edsticas", "Fallos similares",
             "Tribunales", "Jueces"]
    _badges(sl, Inches(2.2), Inches(4.15), tools, gap=Inches(0.08))

    ly = Inches(4.75)
    legend = [
        (ORANGE_BG, "Agente con reglas + LLM"),
        (ORANGE, "Ejecutor principal (GPT-5o)"),
        (SUCCESS_BG, "Respuesta verificada"),
    ]
    lx = Inches(3)
    for cd, ltxt in legend:
        sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, lx, ly + Pt(3), Pt(10), Pt(10))
        sh.fill.solid()
        sh.fill.fore_color.rgb = cd
        sh.line.fill.background()
        _tb(sl, lx + Pt(14), ly, Inches(2), Inches(0.25),
            ltxt, sz=Pt(9), clr=GRAY_500)
        lx += Inches(2.5)

    _foot(sl, 5)


def s05_techstack(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(0.8), "Stack tecnol\u00f3gico")
    _h(sl, Inches(1.1), "Enterprise-grade, 100% Azure, escalable", sz=Pt(26))

    sections = [
        [("Backend",
          [".NET 8", "ASP.NET Core", "MediatR (CQRS)", "EF Core",
           "FluentValidation", "AutoMapper"], False),
         ("Frontend",
          ["Angular 18", "Standalone Components", "SSE Streaming", "PwC AppKit4"], False),
         ("IA / ML",
          ["GPT-5o", "GPT-5o-mini", "text-embedding-3-large",
           "Function Calling", "Structured Output"], True)],
        [("Azure Services",
          ["Azure SQL", "AI Search", "Blob Storage", "Azure OpenAI",
           "Storage Queues", "Container Apps", "App Service", "Static Web Apps"], False),
         ("Workers (4 microservicios)",
          ["Crawler (Selenium)", "Parser (PdfPig)",
           "Enrichment (GPT-5o)", "Indexer (AI Search)"], False),
         ("DevOps",
          ["GitHub Actions CI/CD", "Docker", "ACR", "Staging/Prod slots"], False)],
    ]

    col_w = Inches(5.5)
    col_gap = Inches(0.5)
    x_cols = [Inches(1.1), Inches(1.1) + col_w + col_gap]
    max_badge_w = col_w - Inches(0.4)

    for ci, col_secs in enumerate(sections):
        y = Inches(2.2)
        for title, tags, is_ai in col_secs:
            ch = Inches(1.15)
            _rr(sl, x_cols[ci], y, col_w, ch)
            _acc(sl, x_cols[ci], y + Pt(3), ch - Pt(6))
            _tb(sl, x_cols[ci] + Inches(0.2), y + Inches(0.1),
                col_w - Inches(0.4), Inches(0.25),
                title, sz=Pt(12), clr=GRAY_900, b=True)
            _badges(sl, x_cols[ci] + Inches(0.2), y + Inches(0.42), tags,
                    bg=ORANGE_BG if is_ai else LIGHT_GRAY_BG,
                    fg=ORANGE if is_ai else LIGHT_GRAY_FG,
                    gap=Inches(0.06), max_w=max_badge_w, fsz=Pt(8))
            y += ch + Inches(0.15)

    _foot(sl, 6)


def s06_sources(prs):
    sl = _sl(prs, HERO_BG)
    _eye(sl, Inches(0.8), "Cobertura de fuentes")
    _h(sl, Inches(1.1),
       "Ingesta multi-fuente de jurisprudencia argentina", sz=Pt(26))

    sources = [
        ("CSJN", "Corte Suprema de Justicia\nde la Naci\u00f3n",
         "Phase 1 \u2713", True, "API-first + Selenium\nGPT-5o solo para gaps"),
        ("SAIJ", "Sistema Argentino de\nInformaci\u00f3n Jur\u00eddica",
         "Phase 2", False, "HTML + PDF\nEnrichment completo GPT-5o"),
        ("PJN", "Poder Judicial\nde la Naci\u00f3n",
         "Phase 2", False, "HTML + PDF\nMetadata b\u00e1sica + enrichment"),
        ("SCBA", "Suprema Corte\nde Buenos Aires",
         "Phase 2", False, "HTML + PDF\nMetadata b\u00e1sica + enrichment"),
    ]

    cw, ch = Inches(2.7), Inches(3.2)
    gap = Inches(0.3)
    x0 = (SW - 4 * cw - 3 * gap) // 2
    y0 = Inches(2.3)

    for i, (code, name, phase, done, detail) in enumerate(sources):
        x = x0 + i * (cw + gap)
        _rr(sl, x, y0, cw, ch)
        _acc(sl, x, y0 + Pt(3), ch - Pt(6))
        clr_c = ORANGE if done else GRAY_500
        _tb(sl, x, y0 + Inches(0.25), cw, Inches(0.4),
            code, fnt=HD, sz=Pt(28), clr=clr_c, al=PP_ALIGN.CENTER)
        _tb(sl, x + Inches(0.15), y0 + Inches(0.75), cw - Inches(0.3), Inches(0.5),
            name, sz=Pt(11), clr=GRAY_800, al=PP_ALIGN.CENTER)

        bw = Pt(len(phase) * 5.8 + 24)
        bx = x + (cw - bw) // 2
        bg_b = SUCCESS_BG if done else ORANGE_BG
        fg_b = SUCCESS if done else ORANGE
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  bx, y0 + Inches(1.45), bw, Pt(22))
        sh.fill.solid()
        sh.fill.fore_color.rgb = bg_b
        sh.line.fill.background()
        sh.adjustments[0] = 0.5
        tf = sh.text_frame
        tf.word_wrap = False
        tf.margin_top = Pt(1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = phase
        r.font.name = BD
        r.font.size = Pt(9)
        r.font.color.rgb = fg_b
        r.font.bold = True

        _tb(sl, x + Inches(0.15), y0 + Inches(1.85), cw - Inches(0.3), Inches(0.6),
            detail, sz=Pt(10), clr=GRAY_500, al=PP_ALIGN.CENTER)

    _tb(sl, Inches(1.5), Inches(5.8), Inches(10.333), Inches(0.4),
        "Cada fallo es enriquecido con: jueces firmantes, leyes citadas, tipo de cita "
        "(UPHOLDS, OVERRULES, DISTINGUISHES, CITES), keywords y resumen + holding.",
        sz=Pt(12), clr=GRAY_500, al=PP_ALIGN.CENTER)
    _foot(sl, 7)


def s07_ontology(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(0.5), "Arquitectura de conocimiento")
    _h(sl, Inches(0.8),
       "De la documentaci\u00f3n dispersa al conocimiento ontol\u00f3gico", sz=Pt(24))
    _tb(sl, Inches(2), Inches(1.5), Inches(9.333), Inches(0.4),
        "La informaci\u00f3n jur\u00eddica fragmentada se transforma en una base de "
        "conocimiento estructurada sem\u00e1nticamente, que potencia agentes IA especializados.",
        sz=Pt(11), clr=GRAY_500, al=PP_ALIGN.CENTER)

    BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
    BLUE_FG = RGBColor(0x2B, 0x7B, 0xB9)

    col_y = Inches(2.15)

    # --- Column 1: External Legal Documentation ---
    c1_x, c1_w = Inches(0.45), Inches(2.8)
    _tb(sl, c1_x, col_y, c1_w, Inches(0.25),
        "DOCUMENTACI\u00d3N LEGAL EXTERNA", sz=Pt(8), clr=GRAY_500,
        b=True, al=PP_ALIGN.CENTER)

    sources = [
        ("CSJN", "Corte Suprema"),
        ("SAIJ", "Info Jur\u00eddica"),
        ("PJN", "Poder Judicial"),
        ("SCBA", "Sup. Corte BsAs"),
    ]
    box_w, box_h = Inches(1.25), Inches(0.85)
    box_gap = Inches(0.15)
    bx0 = c1_x + (c1_w - 2 * box_w - box_gap) // 2
    by0 = col_y + Inches(0.35)

    for i, (code, name) in enumerate(sources):
        row, col = divmod(i, 2)
        bx = bx0 + col * (box_w + box_gap)
        by = by0 + row * (box_h + box_gap)
        _rr(sl, bx, by, box_w, box_h, fill=WHITE, line=GRAY_300)
        circ = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                    bx + (box_w - Pt(24)) // 2,
                                    by + Inches(0.1), Pt(24), Pt(24))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE_BG
        circ.line.fill.background()
        _tb(sl, bx, by + Inches(0.08), box_w, Pt(26),
            "\U0001f4dc", sz=Pt(10), al=PP_ALIGN.CENTER, clr=BLUE_FG)
        _tb(sl, bx, by + Inches(0.38), box_w, Inches(0.2),
            code, sz=Pt(10), clr=GRAY_900, b=True, al=PP_ALIGN.CENTER)
        _tb(sl, bx, by + Inches(0.55), box_w, Inches(0.2),
            name, sz=Pt(7), clr=GRAY_500, al=PP_ALIGN.CENTER)

    _tb(sl, c1_x, by0 + 2 * box_h + box_gap + Inches(0.12), c1_w, Inches(0.3),
        "Fallos, resoluciones, acordadas,\ndictámenes, legislaci\u00f3n",
        sz=Pt(8), clr=GRAY_500, al=PP_ALIGN.CENTER, it=True)

    # --- Arrow 1 ---
    a1_x = c1_x + c1_w + Inches(0.05)
    a_w = Inches(0.8)
    a_mid_y = col_y + Inches(1.2)
    _tb(sl, a1_x, a_mid_y - Pt(6), a_w, Inches(0.3),
        "\u2192", sz=Pt(24), clr=ORANGE, al=PP_ALIGN.CENTER)
    _tb(sl, a1_x, a_mid_y + Inches(0.15), a_w, Inches(0.25),
        "Crawling +\nParsing", sz=Pt(7), clr=ORANGE, b=True, al=PP_ALIGN.CENTER)

    # --- Column 2: Legal Ontology ---
    c2_x = a1_x + a_w + Inches(0.05)
    c2_w = Inches(2.5)
    _tb(sl, c2_x, col_y, c2_w, Inches(0.25),
        "ONTOLOG\u00cdA LEGAL", sz=Pt(8), clr=ORANGE, b=True, al=PP_ALIGN.CENTER)

    ont_h = Inches(2.2)
    _rr(sl, c2_x, col_y + Inches(0.35), c2_w, ont_h,
        fill=ORANGE_BG, line=ORANGE, lw=Pt(1.5))

    circ2 = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                 c2_x + (c2_w - Pt(36)) // 2,
                                 col_y + Inches(0.55), Pt(36), Pt(36))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = ORANGE
    circ2.line.fill.background()
    _tb(sl, c2_x + (c2_w - Pt(36)) // 2, col_y + Inches(0.53),
        Pt(36), Pt(36), "\U0001f578", sz=Pt(16), al=PP_ALIGN.CENTER, clr=WHITE)

    _tb(sl, c2_x + Inches(0.15), col_y + Inches(1.05),
        c2_w - Inches(0.3), Inches(0.2),
        "Modelo sem\u00e1ntico", sz=Pt(11), clr=GRAY_900, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, c2_x + Inches(0.15), col_y + Inches(1.3),
        c2_w - Inches(0.3), Inches(0.5),
        "Taxonom\u00eda de conceptos,\nrelaciones y jerarqu\u00edas\ndel derecho argentino",
        sz=Pt(9), clr=GRAY_500, al=PP_ALIGN.CENTER)

    ont_tags = ["Tribunales", "Materias", "Leyes", "Jueces", "Citas", "Voces"]
    _badges(sl, c2_x + Inches(0.15), col_y + Inches(1.9), ont_tags,
            bg=ORANGE_BG, fg=ORANGE, gap=Inches(0.04),
            max_w=c2_w - Inches(0.3), fsz=Pt(7))

    # --- Arrow 2 ---
    a2_x = c2_x + c2_w + Inches(0.05)
    _tb(sl, a2_x, a_mid_y - Pt(6), a_w, Inches(0.3),
        "\u2192", sz=Pt(24), clr=ORANGE, al=PP_ALIGN.CENTER)
    _tb(sl, a2_x, a_mid_y + Inches(0.15), a_w, Inches(0.25),
        "Enrichment +\nIndexing", sz=Pt(7), clr=ORANGE, b=True, al=PP_ALIGN.CENTER)

    # --- Column 3: Ontological Knowledge Base ---
    c3_x = a2_x + a_w + Inches(0.05)
    c3_w = Inches(2.5)
    SUCCESS_BORDER = RGBColor(0x6A, 0xBF, 0x72)
    _tb(sl, c3_x, col_y, c3_w, Inches(0.25),
        "KNOWLEDGE BASE ONTOL\u00d3GICA", sz=Pt(8), clr=SUCCESS,
        b=True, al=PP_ALIGN.CENTER)

    kb_h = Inches(2.2)
    _rr(sl, c3_x, col_y + Inches(0.35), c3_w, kb_h,
        fill=SUCCESS_BG, line=SUCCESS_BORDER, lw=Pt(1.5))

    circ3 = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                 c3_x + (c3_w - Pt(36)) // 2,
                                 col_y + Inches(0.55), Pt(36), Pt(36))
    circ3.fill.solid()
    circ3.fill.fore_color.rgb = SUCCESS
    circ3.line.fill.background()
    _tb(sl, c3_x + (c3_w - Pt(36)) // 2, col_y + Inches(0.53),
        Pt(36), Pt(36), "\U0001f5c4", sz=Pt(16), al=PP_ALIGN.CENTER, clr=WHITE)

    _tb(sl, c3_x + Inches(0.15), col_y + Inches(1.05),
        c3_w - Inches(0.3), Inches(0.2),
        "Base unificada", sz=Pt(11), clr=GRAY_900, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, c3_x + Inches(0.15), col_y + Inches(1.3),
        c3_w - Inches(0.3), Inches(0.5),
        "Datos estructurados +\nvectores sem\u00e1nticos +\ngrafo de relaciones",
        sz=Pt(9), clr=GRAY_500, al=PP_ALIGN.CENTER)

    kb_tags = ["Azure SQL", "AI Search", "Embeddings", "Grafo citas"]
    _badges(sl, c3_x + Inches(0.15), col_y + Inches(1.9), kb_tags,
            bg=SUCCESS_BG, fg=SUCCESS, gap=Inches(0.04),
            max_w=c3_w - Inches(0.3), fsz=Pt(7))

    # --- Arrow 3 ---
    a3_x = c3_x + c3_w + Inches(0.05)
    a3_w = Inches(0.7)
    _tb(sl, a3_x, a_mid_y - Pt(6), a3_w, Inches(0.3),
        "\u2192", sz=Pt(24), clr=ORANGE, al=PP_ALIGN.CENTER)
    _tb(sl, a3_x, a_mid_y + Inches(0.15), a3_w, Inches(0.25),
        "Function\nCalling", sz=Pt(7), clr=ORANGE, b=True, al=PP_ALIGN.CENTER)

    # --- Column 4: AI Agents ---
    c4_x = a3_x + a3_w
    c4_w = Inches(2.4)
    _tb(sl, c4_x, col_y, c4_w, Inches(0.25),
        "AGENTES IA", sz=Pt(8), clr=ORANGE, b=True, al=PP_ALIGN.CENTER)

    agents = [
        ("\U0001f50d", "B\u00fasqueda sem\u00e1ntica", "Consultas en lenguaje natural"),
        ("\U0001f4ac", "Chat RAG ag\u00e9ntico", "Respuestas con citas verificadas"),
        ("\U0001f517", "An\u00e1lisis de relaciones", "Grafo de citas y precedentes"),
        ("\U0001f4c4", "Generaci\u00f3n de memos", "Documentos con fundamentaci\u00f3n"),
    ]
    agent_h = Inches(0.48)
    agent_gap = Inches(0.1)
    ay = col_y + Inches(0.35)

    for ic, title, desc in agents:
        _rr(sl, c4_x, ay, c4_w, agent_h, fill=WHITE, line=ORANGE, lw=Pt(1))
        circ_a = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                      c4_x + Inches(0.12),
                                      ay + (agent_h - Pt(22)) // 2,
                                      Pt(22), Pt(22))
        circ_a.fill.solid()
        circ_a.fill.fore_color.rgb = ORANGE
        circ_a.line.fill.background()
        _tb(sl, c4_x + Inches(0.12),
            ay + (agent_h - Pt(22)) // 2 - Pt(1),
            Pt(22), Pt(22), ic, sz=Pt(9), al=PP_ALIGN.CENTER, clr=WHITE)
        _tb(sl, c4_x + Inches(0.5), ay + Inches(0.04),
            c4_w - Inches(0.6), Inches(0.2),
            title, sz=Pt(9), clr=GRAY_900, b=True)
        _tb(sl, c4_x + Inches(0.5), ay + Inches(0.24),
            c4_w - Inches(0.6), Inches(0.18),
            desc, sz=Pt(7), clr=GRAY_500)
        ay += agent_h + agent_gap

    # --- Bottom summary ---
    sum_y = Inches(5.2)
    labels = [
        ("Datos no estructurados",
         "PDFs, HTML, texto libre de m\u00faltiples\nfuentes con formatos heterog\u00e9neos"),
        ("Conocimiento estructurado",
         "Entidades, relaciones y sem\u00e1ntica\nunificada bajo una ontolog\u00eda legal"),
        ("Inteligencia accionable",
         "Agentes IA que razonan sobre la\nontolog\u00eda para dar respuestas precisas"),
    ]
    lw = Inches(3.2)
    l_gap = Inches(0.3)
    lx0 = (SW - 3 * lw - 2 * l_gap) // 2

    for i, (lt, ld) in enumerate(labels):
        lx = lx0 + i * (lw + l_gap)
        _tb(sl, lx, sum_y, lw, Inches(0.2),
            lt, sz=Pt(10), clr=GRAY_900, b=True, al=PP_ALIGN.CENTER)
        _tb(sl, lx, sum_y + Inches(0.22), lw, Inches(0.4),
            ld, sz=Pt(9), clr=GRAY_500, al=PP_ALIGN.CENTER)
        if i < 2:
            _tb(sl, lx + lw, sum_y + Inches(0.05), Inches(0.3), Inches(0.25),
                "\u27a1", sz=Pt(16), clr=ORANGE, al=PP_ALIGN.CENTER)

    _foot(sl, 8)


def s08_roadmap(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(0.7), "Hoja de ruta")
    _h(sl, Inches(1.0), "5 fases de evoluci\u00f3n planificadas", sz=Pt(26))

    phases = [
        (0, "Fundaciones", "Repo, infra Azure, CI,\ndashboard de roadmap", "done"),
        (1, "Producto m\u00ednimo",
         "CSJN crawler, pipeline,\nsearch, chat, auth, deploy", "current"),
        (2, "Expansi\u00f3n", "SAIJ, PJN, SCBA,\nscheduling, roles", "future"),
        (3, "Inteligencia", "Neo4j, grafo visual,\ntimeline, analytics", "future"),
        (4, "Cobertura total", "Legislaci\u00f3n, perfiles\nde jueces, doctrina", "future"),
    ]

    pw = Inches(2.2)
    gap = Inches(0.15)
    x0 = (SW - 5 * pw - 4 * gap) // 2
    y0 = Inches(2.2)

    for i, (num, title, desc, status) in enumerate(phases):
        x = x0 + i * (pw + gap)
        cx = x + (pw - Pt(28)) // 2
        circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx, y0, Pt(28), Pt(28))
        circ.fill.solid()
        if status == "done":
            circ.fill.fore_color.rgb = SUCCESS_BG
            nc = SUCCESS
        elif status == "current":
            circ.fill.fore_color.rgb = ORANGE
            nc = WHITE
        else:
            circ.fill.fore_color.rgb = GRAY_100
            nc = GRAY_500
        circ.line.fill.background()
        _tb(sl, cx - Pt(2), y0 - Pt(1), Pt(32), Pt(30),
            str(num), sz=Pt(11), clr=nc, b=True, al=PP_ALIGN.CENTER)

        _tb(sl, x, y0 + Inches(0.45), pw, Inches(0.25),
            title, sz=Pt(12), clr=GRAY_900, b=True, al=PP_ALIGN.CENTER)
        _tb(sl, x, y0 + Inches(0.72), pw, Inches(0.5),
            desc, sz=Pt(10), clr=GRAY_800, al=PP_ALIGN.CENTER)

        if status in ("done", "current"):
            bt = "Completada" if status == "done" else "En progreso"
            bw = Pt(len(bt) * 5.8 + 24)
            bx = x + (pw - bw) // 2
            bg_v = SUCCESS_BG if status == "done" else ORANGE_BG
            fg_v = SUCCESS if status == "done" else ORANGE
            sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      bx, y0 + Inches(1.25), bw, Pt(20))
            sh.fill.solid()
            sh.fill.fore_color.rgb = bg_v
            sh.line.fill.background()
            sh.adjustments[0] = 0.5
            tf = sh.text_frame
            tf.word_wrap = False
            tf.margin_top = Pt(1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = bt
            r.font.name = BD
            r.font.size = Pt(8)
            r.font.color.rgb = fg_v
            r.font.bold = True

        if i < 4:
            sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x + pw, y0 + Inches(0.3), Pt(1), Inches(0.8))
            sep.fill.solid()
            sep.fill.fore_color.rgb = GRAY_300
            sep.line.fill.background()

    sw2, sh2 = Inches(3.4), Inches(1.2)
    gap2 = Inches(0.35)
    x0s = (SW - 3 * sw2 - 2 * gap2) // 2
    ys = Inches(4.5)
    stats = [("250", "Entregables planificados"),
             ("33", "Features definidas"),
             ("18", "Features en Phase 1")]
    for i, (v, l) in enumerate(stats):
        _stat(sl, x0s + i * (sw2 + gap2), ys, sw2, sh2, v, l)
    _foot(sl, 9)


def s09_vision(prs):
    sl = _sl(prs, HERO_BG)
    _eye(sl, Inches(0.6), "Visi\u00f3n futura")
    _h(sl, Inches(0.9),
       "M\u00e1s all\u00e1 del MVP: un ecosistema legal inteligente", sz=Pt(26))

    items = [
        ("\U0001f4cb", "Generaci\u00f3n de memos",
         "Memo jur\u00eddico autom\u00e1tico con citas verificadas a partir de una pregunta."),
        ("\u270f", "Asistente de redacci\u00f3n",
         "Drafting assistant para escritos legales con sugerencias de jurisprudencia."),
        ("\U0001f4ca", "Predicci\u00f3n de resultados",
         "An\u00e1lisis predictivo basado en historial de fallos por tema, tribunal y juez."),
        ("\U0001f517", "Grafo visual interactivo",
         "Exploraci\u00f3n visual de relaciones entre fallos, jueces y leyes con Neo4j."),
        ("\U0001f514", "Alertas inteligentes",
         "Notificaciones cuando se indexan fallos relevantes para tus suscripciones."),
        ("\U0001f4c4", "Plugin para Word",
         "B\u00fasqueda jurisprudencial y citaci\u00f3n directa desde Microsoft Word."),
    ]

    cw, ch = Inches(3.4), Inches(1.6)
    gap = Inches(0.35)
    cols = 3
    x0 = (SW - cols * cw - (cols - 1) * gap) // 2

    for i, (ic, ti, bo) in enumerate(items):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gap)
        y = Inches(2.1) + row * (ch + Inches(0.2))
        _card(sl, x, y, cw, ch, icon=ic, title=ti, body=bo)

    _badges(sl, Inches(1.2), Inches(5.9),
            ["Comparaci\u00f3n de documentos", "Carpetas de investigaci\u00f3n",
             "Bookmarks", "Extracci\u00f3n tabular", "API p\u00fablica",
             "Workspace de equipo", "Perfiles de jueces",
             "Tendencias jurisprudenciales"],
            bg=LIGHT_GRAY_BG, fg=LIGHT_GRAY_FG, gap=Inches(0.06))
    _foot(sl, 10)


def s10_value(prs):
    sl = _sl(prs, GRAY_900)
    _eye(sl, Inches(1.0), "Propuesta de valor", clr=AMBER)
    _h(sl, Inches(1.4), "Impacto directo en la pr\u00e1ctica legal",
       sz=Pt(28), clr=WHITE)

    metrics = [
        ("80%", "Reducci\u00f3n en tiempo\nde investigaci\u00f3n"),
        ("100%", "Citas trazables\ny verificadas"),
        ("4", "Fuentes p\u00fablicas\nintegradas"),
        ("24/7", "Disponibilidad\ndel asistente"),
    ]
    cw3, ch3 = Inches(2.5), Inches(1.5)
    gap = Inches(0.3)
    x0 = (SW - 4 * cw3 - 3 * gap) // 2
    y0 = Inches(2.8)

    for i, (val, lbl) in enumerate(metrics):
        x = x0 + i * (cw3 + gap)
        _rr(sl, x, y0, cw3, ch3, fill=DARK_CARD, line=DARK_BORDER)
        _tb(sl, x, y0 + Inches(0.2), cw3, Inches(0.5),
            val, fnt=HD, sz=Pt(34), clr=AMBER, al=PP_ALIGN.CENTER)
        _tb(sl, x, y0 + Inches(0.8), cw3, Inches(0.5),
            lbl, sz=Pt(11), clr=RGBColor(0xB0, 0xB0, 0xB0), al=PP_ALIGN.CENTER)

    qy = Inches(4.8)
    _acc(sl, Inches(3.5), qy, Inches(0.9), clr=AMBER)
    _tb(sl, Inches(3.8), qy, Inches(6), Inches(1),
        "\u201cUn abogado que puede encontrar el precedente exacto en segundos "
        "tiene una ventaja competitiva insuperable.\u201d",
        fnt=HD, sz=Pt(18), clr=RGBColor(0xE8, 0xE8, 0xE8), it=True)
    _foot(sl, 11, clr=RGBColor(0x60, 0x60, 0x60))


def s11_differentiators(prs):
    sl = _sl(prs, WHITE)
    _eye(sl, Inches(0.8), "Diferenciadores")
    _h(sl, Inches(1.1), "\u00bfPor qu\u00e9 Legal AI AR?", sz=Pt(28))

    items = [
        ("\U0001f4cd", "Espec\u00edfico para derecho argentino",
         "No es un ChatGPT gen\u00e9rico. Ingesta real de fuentes oficiales argentinas "
         "con metadata estructurada y verificada."),
        ("\U0001f6e1", "Citas verificadas, no alucinaciones",
         "Output guardrail valida cada cita contra la base de datos. "
         "El abogado conf\u00eda porque puede verificar."),
        ("\u26a1", "Enterprise-grade en Azure",
         "Infraestructura PwC: Azure Entra ID (SSO), cumplimiento, "
         "scale-to-zero, CI/CD automatizado."),
        ("\u2699", "Agentic, no solo RAG",
         "El modelo decide qu\u00e9 herramientas usar. Puede buscar, cruzar citas, "
         "consultar leyes y agregar estad\u00edsticas."),
        ("\u27e8/\u27e9", "Dise\u00f1ado para evolucionar",
         "Roadmap de 5 fases con 250 entregables. Arquitectura extensible "
         "para nuevas fuentes y capacidades."),
        ("\U0001f9e9", "UX profesional PwC",
         "Dise\u00f1o alineado con AppKit4 de PwC. UI en espa\u00f1ol. "
         "Accesible, responsive, integrada con la marca."),
    ]

    cw, ch = Inches(5.3), Inches(1.5)
    gap = Inches(0.35)
    x0 = (SW - 2 * cw - gap) // 2
    y0 = Inches(2.2)

    for i, (ic, ti, bo) in enumerate(items):
        row, col = divmod(i, 2)
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + Inches(0.2))
        _card(sl, x, y, cw, ch, icon=ic, title=ti, body=bo)
    _foot(sl, 12)


def s12_close(prs):
    sl = _sl(prs, GRAY_900)
    _tb(sl, Inches(5.3), Inches(1.2), Inches(1), Inches(0.4),
        "pwc", fnt=HD, sz=Pt(26), clr=AMBER, b=True)
    sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(6.4), Inches(1.28), Pt(1), Inches(0.25))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
    sep.line.fill.background()
    _tb(sl, Inches(6.6), Inches(1.22), Inches(2), Inches(0.4),
        "Legal AI AR", sz=Pt(14), clr=WHITE)

    _h(sl, Inches(2.0), "La jurisprudencia argentina,\npotenciada por IA",
       sz=Pt(34), clr=WHITE)
    _tb(sl, Inches(2.5), Inches(3.3), Inches(8.333), Inches(0.6),
        "Transformemos la forma en que nuestros abogados investigan, "
        "analizan y construyen argumentos legales.",
        sz=Pt(14), clr=RGBColor(0x99, 0x99, 0x99), al=PP_ALIGN.CENTER)

    _tb(sl, Inches(5.5), Inches(4.1), Inches(2.333), Inches(0.3),
        "\u00bfPr\u00f3ximos pasos?", sz=Pt(13), clr=AMBER, b=True, al=PP_ALIGN.CENTER)

    ctas = [
        ("\U0001f465", "Demo con el\nequipo legal"),
        ("\U0001f5c4", "Piloto con 100+\nfallos CSJN"),
        ("\u2601", "Deploy en\nstaging Azure"),
        ("\U0001f4ca", "Feedback y\nPhase 2"),
    ]
    cw4, ch4 = Inches(2.3), Inches(1.2)
    gap3 = Inches(0.3)
    x0 = (SW - 4 * cw4 - 3 * gap3) // 2
    y0 = Inches(4.7)

    for i, (ic, txt) in enumerate(ctas):
        x = x0 + i * (cw4 + gap3)
        _rr(sl, x, y0, cw4, ch4, fill=DARK_CARD, line=DARK_BORDER)
        _tb(sl, x, y0 + Inches(0.15), cw4, Inches(0.35),
            ic, sz=Pt(22), clr=AMBER, al=PP_ALIGN.CENTER)
        _tb(sl, x, y0 + Inches(0.55), cw4, Inches(0.5),
            txt, sz=Pt(12), clr=RGBColor(0xCC, 0xCC, 0xCC), al=PP_ALIGN.CENTER)

    _foot(sl, 13, clr=RGBColor(0x50, 0x50, 0x50))


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    s00_cover(prs)
    s01_problem(prs)
    s02_solution(prs)
    s03_architecture(prs)
    s04_agentic(prs)
    s05_techstack(prs)
    s06_sources(prs)
    s07_ontology(prs)
    s08_roadmap(prs)
    s09_vision(prs)
    s10_value(prs)
    s11_differentiators(prs)
    s12_close(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "legal-ai-ar-pitch.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
